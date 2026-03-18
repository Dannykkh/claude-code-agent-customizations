# requires: python-pptx Pillow
# pip install python-pptx Pillow

"""
python-pptx 기반 PowerPoint 헬퍼 함수 라이브러리.

각 슬라이드 유형(타이틀, 내용, 이미지, 차트, 테이블)을 독립 함수로 제공하여
SKILL.md의 PPTGenerator 클래스를 사용하지 않고도 빠르게 슬라이드를 조합할 수 있습니다.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE


# ──────────────────────────────────────────────────────────────────────────────
# 프레젠테이션 생성
# ──────────────────────────────────────────────────────────────────────────────

def create_presentation(
    title: str,
    subtitle: str = "",
    author: str = "",
    template_path: str | None = None,
) -> Presentation:
    """
    기본 타이틀 슬라이드가 포함된 새 Presentation 객체를 생성합니다.

    Args:
        title:         발표 제목
        subtitle:      부제목 (선택)
        author:        작성자 이름 (코어 프로퍼티에 기록)
        template_path: 기존 .pptx 템플릿 경로. None이면 기본 테마 사용.

    Returns:
        Presentation 객체 (타이틀 슬라이드 1장 포함)
    """
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    # 코어 프로퍼티 설정
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    prs.core_properties.title = title
    if author:
        prs.core_properties.author = author

    # 타이틀 슬라이드 추가
    add_title_slide(prs, title, subtitle)
    return prs


# ──────────────────────────────────────────────────────────────────────────────
# 타이틀 슬라이드
# ──────────────────────────────────────────────────────────────────────────────

def add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    """
    프레젠테이션에 타이틀 슬라이드를 추가합니다.

    슬라이드 레이아웃 인덱스 0 (제목 슬라이드) 을 사용합니다.

    Args:
        prs:      Presentation 객체
        title:    슬라이드 제목
        subtitle: 부제목 (선택)

    Returns:
        추가된 Slide 객체
    """
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = title
    _apply_title_style(slide.shapes.title, size=Pt(40))

    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
        _apply_body_style(slide.placeholders[1], size=Pt(24))

    return slide


# ──────────────────────────────────────────────────────────────────────────────
# 내용 슬라이드 (불릿 목록)
# ──────────────────────────────────────────────────────────────────────────────

def add_content_slide(prs: Presentation, title: str, bullets: list[str]):
    """
    제목과 불릿 목록이 있는 내용 슬라이드를 추가합니다.

    들여쓰기: 문자열 앞에 '  ' (공백 2개) 를 붙이면 하위 항목으로 처리됩니다.
    예: bullets=["주항목", "  하위항목", "다음항목"]

    Args:
        prs:     Presentation 객체
        title:   슬라이드 제목
        bullets: 불릿 문자열 목록

    Returns:
        추가된 Slide 객체
    """
    layout = prs.slide_layouts[1]  # 제목 + 내용
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    _apply_title_style(slide.shapes.title)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        level = 1 if bullet.startswith("  ") else 0
        text = bullet.strip()
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level
        _apply_paragraph_style(p, size=Pt(18) if level == 0 else Pt(16))

    return slide


# ──────────────────────────────────────────────────────────────────────────────
# 이미지 슬라이드
# ──────────────────────────────────────────────────────────────────────────────

def add_image_slide(
    prs: Presentation,
    title: str,
    image_path: str,
    caption: str = "",
):
    """
    제목과 이미지(+ 선택적 캡션)가 있는 슬라이드를 추가합니다.

    이미지는 슬라이드 중앙(좌: 1.5인치, 상: 2인치, 너비: 7인치)에 배치됩니다.

    Args:
        prs:        Presentation 객체
        title:      슬라이드 제목
        image_path: 이미지 파일 경로 (.png, .jpg 등)
        caption:    이미지 아래 캡션 (선택)

    Returns:
        추가된 Slide 객체
    """
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"이미지 파일을 찾을 수 없습니다: {image_path}")

    layout = prs.slide_layouts[5]  # 제목만
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    _apply_title_style(slide.shapes.title)

    slide.shapes.add_picture(
        image_path,
        Inches(1.5), Inches(1.8),
        width=Inches(7),
    )

    if caption:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = caption
        p.alignment = PP_ALIGN.CENTER
        _apply_paragraph_style(p, size=Pt(14))

    return slide


# ──────────────────────────────────────────────────────────────────────────────
# 차트 슬라이드
# ──────────────────────────────────────────────────────────────────────────────

def add_chart_slide(
    prs: Presentation,
    title: str,
    chart_data: dict,
    chart_type: str = "bar",
):
    """
    제목과 차트가 있는 슬라이드를 추가합니다.

    chart_data 형식:
    {
        "categories": ["Q1", "Q2", "Q3"],
        "series": [
            {"name": "매출", "values": [100, 200, 150]},
            {"name": "비용", "values": [80, 120, 100]},
        ]
    }

    Args:
        prs:        Presentation 객체
        title:      슬라이드 제목
        chart_data: 위 형식의 차트 데이터
        chart_type: "bar" | "column" | "line" | "pie"
                    (기본값: "bar" — 수평 막대)

    Returns:
        추가된 Slide 객체
    """
    _chart_type_map = {
        "bar":    XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line":   XL_CHART_TYPE.LINE,
        "pie":    XL_CHART_TYPE.PIE,
    }
    xl_type = _chart_type_map.get(chart_type.lower(), XL_CHART_TYPE.BAR_CLUSTERED)

    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    _apply_title_style(slide.shapes.title)

    cd = ChartData()
    cd.categories = chart_data.get("categories", [])
    for series in chart_data.get("series", []):
        cd.add_series(series["name"], series["values"])

    chart_frame = slide.shapes.add_chart(
        xl_type, Inches(1), Inches(1.5), Inches(8), Inches(5), cd
    )
    chart = chart_frame.chart
    chart.has_legend = len(chart_data.get("series", [])) > 1

    return slide


# ──────────────────────────────────────────────────────────────────────────────
# 테이블 슬라이드
# ──────────────────────────────────────────────────────────────────────────────

def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list],
):
    """
    제목과 데이터 테이블이 있는 슬라이드를 추가합니다.

    헤더 행은 진한 파란 배경 + 흰 글자로 스타일링됩니다.
    데이터 행은 교번 색상(흰색 / 연한 파란색)으로 가독성을 높입니다.

    Args:
        prs:     Presentation 객체
        title:   슬라이드 제목
        headers: 헤더 문자열 목록
        rows:    데이터 행 목록 (각 행은 headers와 같은 길이)

    Returns:
        추가된 Slide 객체
    """
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    _apply_title_style(slide.shapes.title)

    col_count = len(headers)
    row_count = len(rows) + 1  # 헤더 포함

    table_shape = slide.shapes.add_table(
        row_count, col_count,
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(5),
    )
    table = table_shape.table

    # 열 너비 균등 분배
    col_width = Emu(int(Inches(9) / col_count))
    for i in range(col_count):
        table.columns[i].width = col_width

    # 헤더 행
    header_bg = RGBColor(0x1F, 0x4E, 0x79)
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = str(text)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.bold = True
        run.font.size = Pt(14)

    # 데이터 행
    even_bg = RGBColor(0xD6, 0xE4, 0xF0)
    for row_idx, row_data in enumerate(rows):
        for col_idx, value in enumerate(row_data[:col_count]):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value) if value is not None else ""
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = even_bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            if p.runs:
                p.runs[0].font.size = Pt(13)

    return slide


# ──────────────────────────────────────────────────────────────────────────────
# 내부 스타일 헬퍼
# ──────────────────────────────────────────────────────────────────────────────

def _apply_title_style(shape, size: Emu = Pt(32)) -> None:
    """제목 플레이스홀더에 기본 스타일을 적용하는 내부 헬퍼."""
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = size
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)


def _apply_body_style(shape, size: Emu = Pt(20)) -> None:
    """본문 플레이스홀더에 기본 스타일을 적용하는 내부 헬퍼."""
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = size


def _apply_paragraph_style(para, size: Emu = Pt(18)) -> None:
    """단락에 폰트 크기를 설정하는 내부 헬퍼."""
    for run in para.runs:
        run.font.size = size


# ──────────────────────────────────────────────────────────────────────────────
# 사용 예시 (직접 실행 시)
# ──────────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    prs = create_presentation("2025 연간 보고서", subtitle="마케팅팀", author="홍길동")

    add_content_slide(prs, "주요 성과", [
        "매출 전년 대비 +23% 달성",
        "  Q3 최고 실적 (분기 기준)",
        "신규 고객 1,200명 확보",
        "고객 만족도 92점 달성",
    ])

    add_chart_slide(prs, "분기별 매출 현황", {
        "categories": ["Q1", "Q2", "Q3", "Q4"],
        "series": [
            {"name": "2024", "values": [120, 140, 180, 200]},
            {"name": "2025", "values": [150, 170, 220, 240]},
        ],
    }, chart_type="column")

    add_table_slide(prs, "팀별 실적 요약",
        headers=["팀", "목표(억)", "실적(억)", "달성률"],
        rows=[
            ["영업1팀", 100, 115, "115%"],
            ["영업2팀", 80, 72, "90%"],
            ["마케팅",  50, 55, "110%"],
        ],
    )

    prs.save("sample_presentation.pptx")
    print("PPT 저장 완료: sample_presentation.pptx")
